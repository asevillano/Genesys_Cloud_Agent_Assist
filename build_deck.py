"""Generate the Agent Assist IP deck by cloning the Voice Live Connector deck
and replacing the textual content. Preserves all theme, shapes and visuals.
"""
from pathlib import Path
from copy import copy
import shutil
from pptx import Presentation

REF = Path(r"C:\Users\angels\AppData\Local\Temp\ref.pptx")
OUT = Path(r"c:\Angel\AI GBB\Demos\Genesys_Cloud_Agent_Assist\Genesys Cloud Agent Assist.pptx")

shutil.copy(REF, OUT)
p = Presentation(OUT)


def set_text(shape, new_text: str):
    """Replace a shape's text while keeping the formatting of its first run."""
    tf = shape.text_frame
    # Reuse the first paragraph; remove the rest.
    paragraphs = list(tf.paragraphs)
    first = paragraphs[0]
    # Drop extra paragraphs
    for extra in paragraphs[1:]:
        extra._p.getparent().remove(extra._p)

    runs = list(first.runs)
    if not runs:
        first.text = new_text
        return
    # Keep only the first run and assign new text
    for extra in runs[1:]:
        extra._r.getparent().remove(extra._r)
    runs[0].text = new_text


def set_paragraphs(shape, items):
    """Replace text frame with multiple paragraphs/lines. items is a list of
    strings or (text, bold) tuples. Format of first existing run is reused."""
    tf = shape.text_frame
    paragraphs = list(tf.paragraphs)
    # Capture formatting from the first run if available
    first_run = None
    for para in paragraphs:
        if para.runs:
            first_run = para.runs[0]
            break

    # Wipe everything
    for para in paragraphs:
        para._p.getparent().remove(para._p)

    txBody = tf._txBody
    from pptx.oxml.ns import qn
    from lxml import etree

    for idx, item in enumerate(items):
        text, bold = (item, None) if isinstance(item, str) else item
        p_el = etree.SubElement(txBody, qn("a:p"))
        if first_run is not None:
            # Copy paragraph properties from original first paragraph if any
            orig_pPr = paragraphs[0]._p.find(qn("a:pPr")) if paragraphs else None
            if orig_pPr is not None:
                p_el.insert(0, copy(orig_pPr))
        r_el = etree.SubElement(p_el, qn("a:r"))
        if first_run is not None:
            rPr = first_run._r.find(qn("a:rPr"))
            if rPr is not None:
                new_rPr = copy(rPr)
                if bold is True:
                    new_rPr.set("b", "1")
                elif bold is False:
                    new_rPr.set("b", "0")
                r_el.append(new_rPr)
        t_el = etree.SubElement(r_el, qn("a:t"))
        t_el.text = text


# ============================================================
# SLIDE 1 — Title
# ============================================================
s1 = p.slides[0]
for sh in s1.shapes:
    if not sh.has_text_frame:
        continue
    name = sh.name
    if name == "Rectangle: Rounded Corners 2":
        set_text(sh, "Genesys Cloud Agent Assist")
    elif name == "TextBox 3":
        set_text(
            sh,
            "Real-time transcription, contextual suggestions and wrap-up summaries "
            "for Genesys Cloud agents — powered by Azure OpenAI Realtime and Azure AI Foundry",
        )
    elif name == "TextBox 4":
        set_text(sh, "Who to call: Angel Sevillano")
    elif name == "Rounded Rectangle 20":
        set_text(sh, "https://github.com/asevillano/Genesys_Cloud_Agent_Assist")


# ============================================================
# SLIDE 2 — Overview
# ============================================================
s2 = p.slides[1]


def find_in_group(group, target_name):
    for sh in group.shapes:
        if sh.name == target_name:
            return sh
        if sh.shape_type == 6:
            r = find_in_group(sh, target_name)
            if r is not None:
                return r
    return None


for sh in s2.shapes:
    if sh.shape_type == 6:  # group
        # Group 16 (customer stories) and Group 13 (industry/personas)
        if sh.name == "Group 16":
            sub_title = find_in_group(sh, "Rounded Rectangle 20")
            if sub_title and sub_title.has_text_frame:
                set_text(sub_title, "Customer stories")
            body = find_in_group(sh, "Rectangle 14")
            if body and body.has_text_frame:
                set_text(
                    body,
                    "Engaged with a Tier-1 Spanish Telco (S500) for the Contact Center "
                    "of 3M+ customers across multiple geographies, competing against "
                    "Google and Amazon. Agent Assist complements the Voice Live "
                    "Connector by adding live transcription, KB-grounded suggestions "
                    "and automatic call wrap-up to the agent desktop.",
                )
        elif sh.name == "Group 13":
            tb = find_in_group(sh, "TextBox 10")
            if tb and tb.has_text_frame:
                set_paragraphs(
                    tb,
                    [
                        "Industry",
                        "Telco, Banking, Insurance, Public sector — any voice contact center.",
                        "",
                        "Personas",
                        "CTO, CIO, Head of AI, Head of Contact Center, Quality Assurance lead.",
                        "",
                        "Customer KPIs and Metrics",
                        "AHT reduction; First Contact Resolution; After-Call Work time; CSAT/NPS; QA coverage %.",
                    ],
                )
        continue

    if not sh.has_text_frame:
        continue
    name = sh.name
    if name == "Rounded Rectangle 20":
        set_text(sh, "Overview")
    elif name == "Title 1":
        set_text(sh, "Asset overview")
    elif name == "TextBox 2":
        set_text(sh, "Genesys Cloud Agent Assist")
    elif name == "Rectangle 5":
        set_paragraphs(
            sh,
            [
                "Description",
                "Genesys Cloud Agent Assist streams live call audio from Genesys "
                "AudioHook v2 into two parallel Azure OpenAI Realtime sessions "
                "(customer + agent channels) for dual-channel transcription. Final "
                "customer turns are sent to an Azure AI Foundry agent (Responses "
                "API) that returns grounded suggestions in real time. On wrap-up, an "
                "Azure OpenAI chat model produces a summary and category tagging.",
                "All turns, suggestions and summaries are persisted in Cosmos DB "
                "(partition key /conversationId). The UI is delivered as a Genesys "
                "Premium Client App iframe. A built-in browser simulator emulates "
                "the Genesys agent desktop + AudioHook connector so the full flow "
                "can be demoed without a Genesys license.",
            ],
        )
    elif name == "TextBox 6":
        set_paragraphs(
            sh,
            [
                "Business Impact ",
                "Reduces Average Handle Time by surfacing relevant knowledge to "
                "agents in real time, improves First Contact Resolution with always-"
                "on call context, and standardises After-Call Work via automated "
                "wrap-up and category tagging — freeing 15-30% of agent capacity for "
                "high-value interactions while improving QA coverage to 100%.",
            ],
        )
    elif name == "TextBox 7":
        set_paragraphs(
            sh,
            [
                "Location and Usage",
                "GitHub repo: https://github.com/asevillano/Genesys_Cloud_Agent_Assist",
            ],
        )


# ============================================================
# SLIDE 3 — Demo Video
# ============================================================
s3 = p.slides[2]
for sh in s3.shapes:
    if not sh.has_text_frame:
        continue
    if sh.name == "Title 1":
        set_text(sh, "Demo Video")
    elif sh.name == "TextBox 1":
        set_text(sh, "Genesys Cloud Agent Assist")


# ============================================================
# SLIDE 4 — Architecture
# ============================================================
s4 = p.slides[3]
arch_map = {
    "Title 1": "Asset architecture",
    "TextBox 9": "Short description",
    "TextBox 3": "Customer",
    "TextBox 15": "Microsoft Foundry",
    "TextBox 21": "gpt-4o-mini-transcribe + gpt-4.1",
    "TextBox 31": "Realtime STT + Foundry Agent",
    "TextBox 56": "GC-Agent-Assist",
    "TextBox 57": "Genesys Cloud Agent Assist (ACA)",
    "TextBox 92": "Agent Desktop / Simulator",
    "TextBox 104": "WS 8kHz",
    "TextBox 106": "WS 24kHz",
    "TextBox 112": "WS 8kHz",
    "TextBox 113": "AudioHook v2",
    "TextBox 8": "Genesys Cloud",
}
for sh in s4.shapes:
    if not sh.has_text_frame:
        continue
    if sh.name in arch_map:
        set_text(sh, arch_map[sh.name])
    elif sh.name == "TextBox 10":
        set_paragraphs(
            sh,
            [
                "Genesys Cloud uses AudioHook v2 to stream stereo 8 kHz µ-law audio "
                "(customer + agent channels) over WebSockets to the Agent Assist backend.",
                "The Agent Assist backend (FastAPI on Azure Container Apps) "
                "deinterleaves the channels, upsamples each to 24 kHz PCM16 and "
                "feeds them into two independent Azure OpenAI Realtime sessions "
                "(gpt-4o-mini-transcribe).",
                "Final customer turns are forwarded to an Azure AI Foundry agent "
                "(Responses API, name-based reference) which streams grounded "
                "suggestions back to the agent desktop iframe.",
                "On call wrap-up an Azure OpenAI chat model produces a summary "
                "and a list of detected categories; all artefacts are persisted "
                "in Azure Cosmos DB partitioned by conversationId.",
                "The browser-based simulator emulates Genesys Edge + agent desktop "
                "for local testing prior to the real Genesys Cloud integration.",
            ],
        )
    elif sh.name == "TextBox 20":
        set_paragraphs(sh, ["Cosmos DB", "(persistence)"])


p.save(OUT)
print(f"Saved: {OUT}")
